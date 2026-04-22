"""
Script para generar la presentación PPTX - Tutorial de Funcionalidades
del Sistema de Transcripción y Análisis de Audios (T&A Hub).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── THEME COLORS ─────────────────────────────────────────────────────────────

DARK_BG = RGBColor(11, 18, 32)       # #0B1220  - Sidebar dark
ACCENT_CYAN = RGBColor(0, 188, 212)   # #00BCD4  - Cyan accent
ACCENT_BLUE = RGBColor(43, 87, 154)   # #2B579A  - Blue accent
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
DARK_TEXT = RGBColor(30, 30, 30)
MEDIUM_GRAY = RGBColor(100, 100, 100)
SOFT_BG = RGBColor(245, 247, 250)
SUCCESS_GREEN = RGBColor(34, 197, 94)
WARNING_AMBER = RGBColor(245, 158, 11)
ERROR_RED = RGBColor(239, 68, 68)
CARD_BG = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── HELPERS ──────────────────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BG):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or CARD_BG
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=14,
                      color=WHITE, font_name="Segoe UI", line_spacing=1.5, bold_first=False):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_CYAN
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, icon_text="",
             title_color=ACCENT_CYAN, bg_color=RGBColor(20, 30, 50)):
    """Card with title and body text."""
    card = add_shape(slide, left, top, width, height, fill_color=bg_color,
                     line_color=RGBColor(50, 70, 100), line_width=Pt(1))
    card.shadow.inherit = False

    y = top + Inches(0.15)
    if icon_text:
        add_text_box(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.35),
                     icon_text, font_size=12, color=MEDIUM_GRAY, bold=False)
        y += Inches(0.3)

    add_text_box(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    y += Inches(0.45)

    add_multiline_box(slide, left + Inches(0.2), y, width - Inches(0.4),
                      height - (y - top) - Inches(0.15), body_lines,
                      font_size=12, color=LIGHT_GRAY)
    return card


def add_step_card(slide, left, top, width, height, step_num, title, description):
    """Numbered step card."""
    card = add_shape(slide, left, top, width, height, fill_color=RGBColor(20, 30, 50),
                     line_color=ACCENT_CYAN, line_width=Pt(2))

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15),
                                     top + Inches(0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_CYAN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(step_num)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, left + Inches(0.8), top + Inches(0.15), width - Inches(1), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.55), width - Inches(1),
                 height - Inches(0.7), description, font_size=11, color=LIGHT_GRAY)
    return card


def add_sidebar_mockup(slide, left, top, width, height, active_item=""):
    """Draw a simplified sidebar mockup."""
    sidebar = add_shape(slide, left, top, width, height, fill_color=DARK_BG,
                        line_color=RGBColor(50, 60, 80), line_width=Pt(1))

    # Title
    add_text_box(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.3),
                 "Admin Panel", font_size=11, color=WHITE, bold=True)

    items = [
        ("PRINCIPAL", None),
        ("⊞  Inicio", "/dashboard"),
        ("GESTIÓN", None),
        ("📢  Campañas", "/campaigns"),
        ("🔊  Audios", "/audios"),
        ("👤  Centro Análisis", "/analysis"),
        ("💬  Prompts", "/prompts"),
    ]

    y_offset = top + Inches(0.5)
    for label, route in items:
        if route is None:
            add_text_box(slide, left + Inches(0.1), y_offset, width - Inches(0.2), Inches(0.22),
                         label, font_size=7, color=MEDIUM_GRAY, bold=True)
            y_offset += Inches(0.22)
        else:
            is_active = active_item and active_item in route
            bg_c = RGBColor(0, 188, 212) if is_active else DARK_BG
            alpha = 30 if is_active else 0
            if is_active:
                item_bg = add_shape(slide, left + Inches(0.05), y_offset,
                                    width - Inches(0.1), Inches(0.28),
                                    fill_color=RGBColor(0, 60, 80),
                                    line_color=ACCENT_CYAN, line_width=Pt(1))
            txt_color = WHITE if is_active else LIGHT_GRAY
            add_text_box(slide, left + Inches(0.15), y_offset, width - Inches(0.3), Inches(0.28),
                         label, font_size=9, color=txt_color, bold=is_active)
            y_offset += Inches(0.32)

    return sidebar


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 - PORTADA
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

# Decorative top bar
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

# Main title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "T&A Hub", font_size=60, color=ACCENT_CYAN, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Transcription & Analysis Platform", font_size=28, color=WHITE,
             bold=False, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(3.6), Inches(3.333))

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "Guía de Funcionalidades del Sistema", font_size=22, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Sistema de transcripción automatizada y análisis de calidad de audios con IA",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Abril 2026  •  Versión 1.0", font_size=12, color=MEDIUM_GRAY,
             alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 - ÍNDICE
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Contenido de la Presentación", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

sections = [
    ("01", "Visión General", "Qué es T&A Hub y para quién está diseñado"),
    ("02", "Inicio de Sesión", "Autenticación y acceso seguro al sistema"),
    ("03", "Dashboard", "Panel principal con estadísticas y métricas"),
    ("04", "Gestión de Campañas", "Crear, editar y administrar campañas de llamadas"),
    ("05", "Gestión de Audios", "Carga, organización y administración de archivos"),
    ("06", "Centro de Análisis", "Pipeline de transcripción y análisis con IA"),
    ("07", "Prompts y Formatos", "Configuración de plantillas de análisis"),
    ("08", "Exportación de Resultados", "Descarga de reportes en Excel"),
    ("09", "Personalización", "Temas visuales y navegación"),
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.5) + row * Inches(1.8)

    add_card(slide, x, y, Inches(3.8), Inches(1.5), title,
             [desc], icon_text=f"Sección {num}")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 - VISIÓN GENERAL
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "01  Visión General", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.6),
             "¿Qué es T&A Hub?", font_size=20, color=ACCENT_CYAN, bold=True)

add_multiline_box(slide, Inches(0.8), Inches(2.0), Inches(6), Inches(2),
                  ["T&A Hub es una plataforma web full-stack diseñada para:",
                   "",
                   "•  Transcribir audios de llamadas telefónicas de forma automática",
                   "•  Analizar la calidad de las llamadas usando Inteligencia Artificial",
                   "•  Gestionar campañas y audios de forma centralizada",
                   "•  Exportar resultados de análisis en formato Excel",
                   "",
                   "El sistema utiliza servicios de IA como Deepgram, WhisperX y",
                   "OpenAI GPT para procesar y evaluar audios masivamente."],
                  font_size=13, color=LIGHT_GRAY)

# Right side - Architecture summary
add_card(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(1.3),
         "Frontend", ["Angular 21 + TypeScript", "TailwindCSS + DaisyUI", "Señales reactivas (Signals)"],
         icon_text="🖥️  CLIENTE")

add_card(slide, Inches(7.5), Inches(3.0), Inches(5), Inches(1.3),
         "Backend", ["FastAPI (Python)", "SQLAlchemy ORM", "JWT Authentication"],
         icon_text="⚙️  SERVIDOR")

add_card(slide, Inches(7.5), Inches(4.6), Inches(5), Inches(1.3),
         "Servicios IA", ["Deepgram (transcripción cloud)", "WhisperX (transcripción local)", "OpenAI GPT (análisis QA)"],
         icon_text="🤖  INTELIGENCIA ARTIFICIAL")

# Users section
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.5),
             "¿Para quién es?", font_size=18, color=ACCENT_CYAN, bold=True)

add_multiline_box(slide, Inches(0.8), Inches(5.1), Inches(6), Inches(1.5),
                  ["👤  Analistas QA — Evalúan calidad de llamadas",
                   "👥  Supervisores — Revisan reportes consolidados",
                   "🔧  Administradores — Gestionan configuración global"],
                  font_size=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 - LOGIN
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "02  Inicio de Sesión", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Login form mockup
login_card = add_shape(slide, Inches(1.5), Inches(1.8), Inches(4.5), Inches(4.5),
                       fill_color=RGBColor(20, 30, 50), line_color=RGBColor(50, 70, 100),
                       line_width=Pt(1))

add_text_box(slide, Inches(2.2), Inches(2.0), Inches(3), Inches(0.5),
             "🔐  Iniciar Sesión", font_size=20, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Username field
add_shape(slide, Inches(2), Inches(2.8), Inches(3.5), Inches(0.5),
          fill_color=RGBColor(30, 40, 60), line_color=RGBColor(80, 90, 110), line_width=Pt(1))
add_text_box(slide, Inches(2.15), Inches(2.85), Inches(3.2), Inches(0.4),
             "👤  Usuario", font_size=12, color=MEDIUM_GRAY)

# Password field
add_shape(slide, Inches(2), Inches(3.5), Inches(3.5), Inches(0.5),
          fill_color=RGBColor(30, 40, 60), line_color=RGBColor(80, 90, 110), line_width=Pt(1))
add_text_box(slide, Inches(2.15), Inches(3.55), Inches(3.2), Inches(0.4),
             "🔒  Contraseña", font_size=12, color=MEDIUM_GRAY)

# Login button
btn = add_shape(slide, Inches(2), Inches(4.3), Inches(3.5), Inches(0.55),
                fill_color=ACCENT_CYAN)
add_text_box(slide, Inches(2), Inches(4.35), Inches(3.5), Inches(0.45),
             "Iniciar Sesión", font_size=14, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Validation rules
add_text_box(slide, Inches(2), Inches(5.1), Inches(3.5), Inches(0.8),
             "Validaciones: usuario requerido,\ncontraseña mín. 8 caracteres, sin espacios",
             font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Right side - explanation
add_text_box(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5),
             "Flujo de Autenticación", font_size=20, color=ACCENT_CYAN, bold=True)

steps_auth = [
    ("1", "Ingreso de Credenciales",
     "El usuario ingresa su nombre de usuario y contraseña en el formulario de login."),
    ("2", "Validación en Servidor",
     "El backend verifica las credenciales contra la BD (bcrypt hash) y genera un token JWT."),
    ("3", "Almacenamiento de Token",
     "El token JWT se guarda en sessionStorage del navegador (se elimina al cerrar pestaña)."),
    ("4", "Acceso al Sistema",
     "El interceptor HTTP inyecta el token en cada petición. Si expira → redirect a login."),
]

for i, (num, title, desc) in enumerate(steps_auth):
    y = Inches(2.5) + i * Inches(1.15)
    add_step_card(slide, Inches(7), y, Inches(5.5), Inches(1.0), num, title, desc)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 - DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "03  Dashboard — Panel Principal", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "El dashboard muestra un resumen operativo de todos los audios subidos por tu usuario.",
             font_size=14, color=LIGHT_GRAY)

# Stats cards mockup
stat_cards = [
    ("AUDIOS TOTALES", "152", "Todos los audios del usuario", ACCENT_CYAN),
    ("TRANSCRITOS", "128", "Con transcripción guardada", SUCCESS_GREEN),
    ("PENDIENTES", "24", "Sin transcripción aún", WARNING_AMBER),
    ("MINUTOS TOTALES", "847.50", "Duración acumulada", ACCENT_BLUE),
    ("COSTO TOTAL", "12.34", "Costo acumulado del usuario", ERROR_RED),
]

for i, (label, value, desc, color) in enumerate(stat_cards):
    x = Inches(0.5) + i * Inches(2.5)
    card = add_shape(slide, x, Inches(2.0), Inches(2.3), Inches(1.5),
                     fill_color=RGBColor(20, 30, 50), line_color=RGBColor(50, 70, 100),
                     line_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2), Inches(0.3),
                 label, font_size=9, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(2), Inches(0.5),
                 value, font_size=28, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.95), Inches(2), Inches(0.3),
                 desc, font_size=9, color=MEDIUM_GRAY)

# Progress section
add_shape(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(2.8),
          fill_color=RGBColor(20, 30, 50), line_color=RGBColor(50, 70, 100), line_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(3.9), Inches(4), Inches(0.4),
             "Progreso de Transcripción", font_size=16, color=WHITE, bold=True)

# Progress bar
add_shape(slide, Inches(0.7), Inches(4.45), Inches(5.5), Inches(0.3),
          fill_color=RGBColor(40, 50, 70))
add_shape(slide, Inches(0.7), Inches(4.45), Inches(4.63), Inches(0.3),
          fill_color=ACCENT_CYAN)
add_text_box(slide, Inches(4.8), Inches(4.4), Inches(1.5), Inches(0.4),
             "84%", font_size=14, color=WHITE, bold=True)

# Stats below progress
add_card(slide, Inches(0.7), Inches(5.0), Inches(2.5), Inches(1.2),
         "TRANSCRITOS", ["128", "84% del total"],
         title_color=SUCCESS_GREEN, bg_color=RGBColor(25, 40, 35))

add_card(slide, Inches(3.5), Inches(5.0), Inches(2.5), Inches(1.2),
         "PENDIENTES", ["24", "16% del total"],
         title_color=WARNING_AMBER, bg_color=RGBColor(40, 35, 20))

# Cost chart section
add_shape(slide, Inches(6.8), Inches(3.8), Inches(6), Inches(2.8),
          fill_color=RGBColor(20, 30, 50), line_color=RGBColor(50, 70, 100), line_width=Pt(1))
add_text_box(slide, Inches(7), Inches(3.9), Inches(5), Inches(0.4),
             "Costo Total por Campaña", font_size=16, color=WHITE, bold=True)

# Mini bar chart mockup
campaigns = [("Campaña Ventas Q1", 4.2), ("Campaña Soporte", 3.1), ("Campaña Retención", 2.8),
             ("QA Mensual", 1.5), ("Piloto IA", 0.7)]
for i, (name, cost) in enumerate(campaigns):
    y = Inches(4.5) + i * Inches(0.4)
    bar_w = Inches(cost / 4.2 * 3.5)
    add_text_box(slide, Inches(7), y, Inches(2.2), Inches(0.35),
                 name, font_size=8, color=LIGHT_GRAY)
    add_shape(slide, Inches(9.3), y + Inches(0.05), bar_w, Inches(0.22),
              fill_color=ACCENT_CYAN)
    add_text_box(slide, Inches(9.3) + bar_w + Inches(0.1), y, Inches(0.8), Inches(0.35),
                 f"${cost:.2f}", font_size=8, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 - CAMPAÑAS
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "04  Gestión de Campañas", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Las campañas organizan los audios en grupos lógicos (ej: campañas de ventas, soporte, QA mensual).",
             font_size=14, color=LIGHT_GRAY)

# Table mockup
table_bg = add_shape(slide, Inches(0.5), Inches(2.0), Inches(7.5), Inches(4.5),
                     fill_color=RGBColor(20, 30, 50), line_color=RGBColor(50, 70, 100),
                     line_width=Pt(1))

# Header bar
add_shape(slide, Inches(0.5), Inches(2.0), Inches(7.5), Inches(0.5), fill_color=RGBColor(30, 45, 70))
add_text_box(slide, Inches(0.7), Inches(2.05), Inches(1.5), Inches(0.4),
             "Nombre", font_size=10, color=WHITE, bold=True)
add_text_box(slide, Inches(2.5), Inches(2.05), Inches(1.5), Inches(0.4),
             "Audios", font_size=10, color=WHITE, bold=True)
add_text_box(slide, Inches(3.7), Inches(2.05), Inches(1.5), Inches(0.4),
             "Estado", font_size=10, color=WHITE, bold=True)
add_text_box(slide, Inches(5.3), Inches(2.05), Inches(1.5), Inches(0.4),
             "Costo", font_size=10, color=WHITE, bold=True)
add_text_box(slide, Inches(6.7), Inches(2.05), Inches(1.2), Inches(0.4),
             "Acciones", font_size=10, color=WHITE, bold=True)

# Sample rows
rows_data = [
    ("Ventas Q1 2026", "45", "COMPLETED", "$4.20"),
    ("Soporte Técnico", "32", "PARTIAL", "$3.10"),
    ("Retención Abril", "28", "PENDING", "$0.00"),
    ("QA Mensual", "15", "COMPLETED", "$1.50"),
    ("Piloto IA", "8", "EMPTY", "$0.00"),
]

for i, (name, audios, status, cost) in enumerate(rows_data):
    y = Inches(2.6) + i * Inches(0.45)
    status_color = {
        "COMPLETED": SUCCESS_GREEN, "PARTIAL": WARNING_AMBER,
        "PENDING": ACCENT_CYAN, "EMPTY": MEDIUM_GRAY
    }.get(status, MEDIUM_GRAY)

    add_text_box(slide, Inches(0.7), y, Inches(1.8), Inches(0.4),
                 name, font_size=10, color=WHITE)
    add_text_box(slide, Inches(2.5), y, Inches(1), Inches(0.4),
                 audios, font_size=10, color=LIGHT_GRAY)
    add_text_box(slide, Inches(3.7), y, Inches(1.5), Inches(0.4),
                 status, font_size=9, color=status_color, bold=True)
    add_text_box(slide, Inches(5.3), y, Inches(1), Inches(0.4),
                 cost, font_size=10, color=LIGHT_GRAY)
    add_text_box(slide, Inches(6.7), y, Inches(1.2), Inches(0.4),
                 "✏️  🗑️", font_size=10, color=ACCENT_CYAN)

# Search bar mockup
add_shape(slide, Inches(0.5), Inches(5.0), Inches(4), Inches(0.4),
          fill_color=RGBColor(30, 40, 60), line_color=RGBColor(80, 90, 110), line_width=Pt(1))
add_text_box(slide, Inches(0.65), Inches(5.0), Inches(3.8), Inches(0.4),
             "🔍  Buscar campaña...", font_size=10, color=MEDIUM_GRAY)

# Pagination
add_text_box(slide, Inches(5), Inches(5.0), Inches(3), Inches(0.4),
             "◀  Página 1 de 5  ▶", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# Right side - CRUD operations
add_text_box(slide, Inches(8.5), Inches(2.0), Inches(4), Inches(0.4),
             "Operaciones Disponibles", font_size=18, color=ACCENT_CYAN, bold=True)

crud_ops = [
    ("➕ Crear Campaña", "Modal con nombre y descripción.\nValidación de campos requeridos."),
    ("✏️ Editar Campaña", "Modifica nombre y descripción.\nSe actualiza en tiempo real."),
    ("🗑️ Eliminar Campaña", "Confirmación requerida.\nElimina audios y análisis asociados (CASCADE)."),
    ("🔍 Buscar", "Búsqueda por nombre con debounce\nde 300ms para rendimiento."),
    ("📊 Estadísticas", "Total audios, transcritos, pendientes,\ncostos de transcripción y análisis."),
]

for i, (title, desc) in enumerate(crud_ops):
    y = Inches(2.6) + i * Inches(0.95)
    add_card(slide, Inches(8.5), y, Inches(4.3), Inches(0.85), title, [desc],
             bg_color=RGBColor(20, 30, 50))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 - AUDIOS
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "05  Gestión de Audios", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Upload section (left)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "Carga de Archivos", font_size=20, color=ACCENT_CYAN, bold=True)

# Drag and drop zone
dnd_zone = add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.2),
                     fill_color=RGBColor(20, 30, 50),
                     line_color=ACCENT_CYAN, line_width=Pt(2))
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(4), Inches(0.5),
             "📂", font_size=40, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(4), Inches(0.4),
             "Arrastra y suelta archivos aquí", font_size=14, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(4), Inches(0.4),
             "o haz click para seleccionar", font_size=11, color=MEDIUM_GRAY,
             alignment=PP_ALIGN.CENTER)

# Format info
add_multiline_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.5),
                  ["Formatos aceptados:  MP3, WAV, M4A, OGG, OGM, MP4",
                   "Tamaño máximo:  500 MB por archivo",
                   "Carga múltiple:  Sí (varios archivos a la vez)",
                   "Detección automática de duración con FFprobe"],
                  font_size=11, color=LIGHT_GRAY)

# Right side - table features
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
             "Tabla de Audios", font_size=20, color=ACCENT_CYAN, bold=True)

audio_features = [
    ("📋 Listado Paginado", "Vista de todos los audios con nombre, campaña, transcripción, duración y costo."),
    ("🔍 Búsqueda + Filtro", "Búsqueda por nombre con debounce. Filtro dropdown por campaña."),
    ("✏️ Editar Nombre", "Modal para renombrar el archivo de audio."),
    ("🗑️ Eliminar Audio", "Elimina el archivo, registro y análisis asociados."),
    ("📊 Estadísticas", "Resumen: total, transcritos, pendientes, minutos totales, costo."),
]

for i, (title, desc) in enumerate(audio_features):
    y = Inches(2.1) + i * Inches(1.05)
    add_card(slide, Inches(7), y, Inches(5.5), Inches(0.9), title, [desc],
             bg_color=RGBColor(20, 30, 50))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 - CENTRO DE ANÁLISIS (Pipeline)
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "06  Centro de Análisis — Pipeline", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "El Centro de Análisis es el corazón del sistema: permite transcribir, analizar y exportar resultados.",
             font_size=14, color=LIGHT_GRAY)

# Pipeline modes
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(4), Inches(0.4),
             "Modos de Pipeline", font_size=18, color=ACCENT_CYAN, bold=True)

modes = [
    ("🎙️ Solo Transcribir", "Convierte audio a texto con diarización de hablantes.",
     "Deepgram (cloud) o WhisperX (local)"),
    ("🤖 Solo Analizar", "Evalúa transcripciones existentes con IA.",
     "OpenAI GPT-4.1-mini"),
    ("⚡ Ambos (Pipeline)", "Transcribe y luego analiza secuencialmente.",
     "Deepgram/WhisperX → OpenAI"),
]

for i, (title, desc, provider) in enumerate(modes):
    y = Inches(2.6) + i * Inches(1.2)
    card = add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(1.05),
                     fill_color=RGBColor(20, 30, 50),
                     line_color=ACCENT_CYAN, line_width=Pt(1))
    add_text_box(slide, Inches(1), y + Inches(0.08), Inches(5), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1), y + Inches(0.4), Inches(5), Inches(0.25),
                 desc, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, Inches(1), y + Inches(0.7), Inches(5), Inches(0.25),
                 f"Proveedor: {provider}", font_size=10, color=ACCENT_CYAN)

# Right side - Pipeline flow
add_text_box(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(0.4),
             "Flujo del Pipeline Completo", font_size=18, color=ACCENT_CYAN, bold=True)

pipeline_steps = [
    ("1", "Seleccionar", "Elige campaña, prompt,\nformato y proveedor"),
    ("2", "Ejecutar", "Click en 'Ejecutar Pipeline'\nInicia tarea en background"),
    ("3", "Transcribir", "Cada audio se envía al\nproveedor de transcripción"),
    ("4", "Analizar", "OpenAI evalúa cada\ntranscripción con el prompt"),
    ("5", "Monitorear", "Polling cada 3s muestra\nprogreso en tiempo real"),
    ("6", "Resultados", "Tabla de resultados y\nopción de exportar Excel"),
]

for i, (num, title, desc) in enumerate(pipeline_steps):
    row = i // 2
    col = i % 2
    x = Inches(7) + col * Inches(2.9)
    y = Inches(2.7) + row * Inches(1.5)
    add_step_card(slide, x, y, Inches(2.7), Inches(1.3), num, title, desc)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 - CENTRO DE ANÁLISIS (Resultados)
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "06  Centro de Análisis — Resultados", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Results table mockup
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "Tabla de Resultados de Análisis", font_size=18, color=ACCENT_CYAN, bold=True)

# Table header
table_x = Inches(0.5)
add_shape(slide, table_x, Inches(2.0), Inches(12), Inches(0.45), fill_color=RGBColor(30, 45, 70))
headers = [("Audio", 2.5), ("Criterio", 2.5), ("Evaluación", 1.5), ("Justificación", 4), ("Costo", 1.2)]
x = table_x + Inches(0.15)
for h_name, h_width in headers:
    add_text_box(slide, x, Inches(2.05), Inches(h_width), Inches(0.35),
                 h_name, font_size=10, color=WHITE, bold=True)
    x += Inches(h_width)

# Sample result rows
results = [
    ("llamada_001.mp3", "Saludo inicial", "✅ Cumple", "El agente saluda correctamente...", "$0.02"),
    ("llamada_001.mp3", "Verificación datos", "❌ No cumple", "No se verificó identidad...", "$0.02"),
    ("llamada_001.mp3", "Cierre de llamada", "✅ Cumple", "Despedida cordial realizada...", "$0.02"),
    ("llamada_002.mp3", "Saludo inicial", "⚠️ No aplica", "Llamada sin conexión inicial...", "$0.01"),
    ("llamada_002.mp3", "Ofrecimiento", "✅ Cumple", "Se ofreció producto correctamente...", "$0.02"),
]

for i, (audio, criterio, evaluacion, justif, costo) in enumerate(results):
    y = Inches(2.55) + i * Inches(0.4)
    bg_c = RGBColor(20, 30, 50) if i % 2 == 0 else RGBColor(25, 35, 55)
    add_shape(slide, table_x, y, Inches(12), Inches(0.38), fill_color=bg_c)
    x = table_x + Inches(0.15)
    for val, w in [(audio, 2.5), (criterio, 2.5), (evaluacion, 1.5), (justif[:35] + "...", 4), (costo, 1.2)]:
        eval_color = WHITE
        if "Cumple" in val and "No" not in val:
            eval_color = SUCCESS_GREEN
        elif "No cumple" in val:
            eval_color = ERROR_RED
        elif "No aplica" in val:
            eval_color = WARNING_AMBER
        add_text_box(slide, x, y + Inches(0.02), Inches(w), Inches(0.33),
                     val, font_size=9, color=eval_color if val == evaluacion else LIGHT_GRAY)
        x += Inches(w)

# Tabs mockup
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(4), Inches(0.4),
             "Vista con dos pestañas:", font_size=14, color=WHITE, bold=True)

tab1 = add_shape(slide, Inches(0.8), Inches(5.3), Inches(3), Inches(0.45),
                 fill_color=ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(5.32), Inches(3), Inches(0.4),
             "📊  Resultados de Análisis", font_size=11, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

tab2 = add_shape(slide, Inches(4), Inches(5.3), Inches(3), Inches(0.45),
                 fill_color=RGBColor(30, 45, 70))
add_text_box(slide, Inches(4), Inches(5.32), Inches(3), Inches(0.4),
             "📝  Transcripciones", font_size=11, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# Features list
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
             "Paginación cliente (10 por página)  •  Filtro por texto  •  Modal de detalle por audio  •  Descarga Excel",
             font_size=11, color=MEDIUM_GRAY)

# Right side - Transcription format
add_text_box(slide, Inches(8), Inches(4.8), Inches(5), Inches(0.4),
             "Formato de Transcripción", font_size=14, color=ACCENT_CYAN, bold=True)

trans_lines = [
    "00:00.000 - 00:03.500 | SPEAKER_00 | Buenos días,",
    "00:03.500 - 00:07.200 | SPEAKER_01 | Hola, le llamo",
    "00:07.200 - 00:12.800 | SPEAKER_00 | Sí, dígame...",
    "00:12.800 - 00:18.100 | SPEAKER_01 | Le ofrezco...",
]

trans_bg = add_shape(slide, Inches(8), Inches(5.3), Inches(5), Inches(1.5),
                     fill_color=RGBColor(15, 20, 35), line_color=RGBColor(50, 70, 100),
                     line_width=Pt(1))

add_multiline_box(slide, Inches(8.15), Inches(5.4), Inches(4.7), Inches(1.3),
                  trans_lines, font_size=9, color=SUCCESS_GREEN, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 - PROMPTS Y FORMATOS
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "07  Prompts y Formatos de Salida", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Prompts section (left)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "📝 Gestión de Prompts", font_size=20, color=ACCENT_CYAN, bold=True)

add_multiline_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1),
                  ["Los prompts son las instrucciones que se envían a la IA para que",
                   "evalúe cada transcripción. Son completamente personalizables."],
                  font_size=12, color=LIGHT_GRAY)

# Prompt example card
prompt_card = add_shape(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2.8),
                        fill_color=RGBColor(15, 20, 35),
                        line_color=RGBColor(50, 70, 100), line_width=Pt(1))

add_text_box(slide, Inches(1), Inches(3.1), Inches(5), Inches(0.3),
             "Ejemplo de Prompt:", font_size=11, color=WARNING_AMBER, bold=True)

prompt_text = [
    "Eres un analista QA de llamadas telefónicas.",
    "Evalúa la siguiente transcripción según estos criterios:",
    "",
    "1. Saludo inicial: ¿El agente saluda correctamente?",
    "2. Verificación de datos: ¿Se verificó la identidad?",
    "3. Ofrecimiento de producto: ¿Se realizó el pitch?",
    "4. Cierre de llamada: ¿Se despidió correctamente?",
    "",
    "Para cada criterio responde: Cumple / No cumple / No aplica",
    "con justificación detallada.",
]

add_multiline_box(slide, Inches(1), Inches(3.45), Inches(5), Inches(2.2),
                  prompt_text, font_size=9, color=SUCCESS_GREEN, font_name="Consolas")

# Operations
add_multiline_box(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.8),
                  ["Operaciones: Crear, Editar (inline), Eliminar, Activar/Desactivar"],
                  font_size=11, color=LIGHT_GRAY)

# Formats section (right)
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
             "📋 Formatos de Salida", font_size=20, color=ACCENT_CYAN, bold=True)

add_multiline_box(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(1),
                  ["Los formatos definen qué campos tendrá el resultado del",
                   "análisis. Configuran las columnas del reporte exportado."],
                  font_size=12, color=LIGHT_GRAY)

# Format example
format_card = add_shape(slide, Inches(7), Inches(3.0), Inches(5.5), Inches(2.3),
                        fill_color=RGBColor(15, 20, 35),
                        line_color=RGBColor(50, 70, 100), line_width=Pt(1))

add_text_box(slide, Inches(7.2), Inches(3.1), Inches(5), Inches(0.3),
             "Ejemplo de Formato:", font_size=11, color=WARNING_AMBER, bold=True)

format_fields = [
    "Nombre: \"Evaluación QA Estándar\"",
    "",
    "Campos definidos:",
    "  • criterio       → Nombre del criterio evaluado",
    "  • evaluacion     → Cumple / No cumple / No aplica",
    "  • justificacion  → Explicación detallada",
    "  • obs_adicional  → Observaciones complementarias",
]

add_multiline_box(slide, Inches(7.2), Inches(3.45), Inches(5), Inches(1.7),
                  format_fields, font_size=10, color=LIGHT_GRAY, font_name="Consolas")

# Suggestions
add_card(slide, Inches(7), Inches(5.5), Inches(5.5), Inches(1.2),
         "💡 Sugerencias Automáticas",
         ["El sistema sugiere combinaciones óptimas de prompt + formato",
          "basándose en las plantillas existentes del usuario.",
          "Incluye puntuación de compatibilidad y razón de la sugerencia."],
         title_color=WARNING_AMBER, bg_color=RGBColor(35, 30, 20))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 - EXPORTACIÓN
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "08  Exportación de Resultados", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Los resultados de análisis se pueden exportar en formato Excel (.xlsx) con estructura profesional.",
             font_size=14, color=LIGHT_GRAY)

# Excel structure mockup
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(6), Inches(0.4),
             "📊 Estructura del Archivo Excel", font_size=18, color=ACCENT_CYAN, bold=True)

# Sheet 1 - Resultados
add_shape(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.4), fill_color=SUCCESS_GREEN)
add_text_box(slide, Inches(1), Inches(2.72), Inches(4), Inches(0.35),
             "Hoja 1: RESULTADOS", font_size=12, color=WHITE, bold=True)

# Excel columns header
cols_header = add_shape(slide, Inches(0.8), Inches(3.15), Inches(11.5), Inches(0.4),
                        fill_color=RGBColor(30, 45, 70))
excel_cols = ["AUDIO_ID", "AUDIO_NOMBRE", "DNI", "EJECUTIVO", "FECHA", "[1] Criterio", "[2] Criterio", "OBS", "TRANSCRIPCIÓN"]
x = Inches(0.9)
for col in excel_cols:
    w = Inches(1.2) if "Criterio" not in col else Inches(1.4)
    add_text_box(slide, x, Inches(3.17), Inches(1.3), Inches(0.35),
                 col, font_size=7, color=WHITE, bold=True)
    x += Inches(1.28)

# Sample data row
data_row = add_shape(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.35),
                     fill_color=RGBColor(240, 245, 250))
sample_data = ["uuid-001", "llamada_001", "12345678", "J. Pérez", "2026-04-15", "Cumple ✅", "No cumple ❌", "Obs...", "00:00-..."]
x = Inches(0.9)
for val in sample_data:
    add_text_box(slide, x, Inches(3.62), Inches(1.3), Inches(0.3),
                 val, font_size=7, color=DARK_TEXT)
    x += Inches(1.28)

# Sheet 2 - Costos
add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.4), fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(4.32), Inches(4), Inches(0.35),
             "Hoja 2: COSTOS", font_size=12, color=WHITE, bold=True)

add_multiline_box(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(1),
                  ["Detalle de costos de transcripción y análisis por audio:",
                   "• Costo de transcripción (Deepgram / WhisperX)",
                   "• Costo de análisis (tokens de entrada/salida OpenAI)",
                   "• Totales acumulados por campaña"],
                  font_size=11, color=LIGHT_GRAY)

# Right side - features
add_text_box(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(0.4),
             "Características de la Exportación", font_size=16, color=ACCENT_CYAN, bold=True)

export_features = [
    "✅  Columnas dinámicas según criterios evaluados",
    "✅  Metadatos extraídos del nombre del archivo",
    "✅  Observaciones agrupadas por rangos de criterios",
    "✅  Transcripción dividida en chunks (32KB)",
    "✅  Formato ancho: 1 fila por audio",
    "✅  Descarga directa desde el navegador",
]

add_multiline_box(slide, Inches(7), Inches(4.8), Inches(5.5), Inches(2),
                  export_features, font_size=12, color=LIGHT_GRAY, line_spacing=1.8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 - PERSONALIZACIÓN (Temas + Sidebar)
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "09  Personalización y Navegación", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# Theme section
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "🎨 Temas Visuales", font_size=20, color=ACCENT_CYAN, bold=True)

# Light theme card
add_card(slide, Inches(0.8), Inches(2.1), Inches(2.8), Inches(2.5),
         "☀️  Tema Claro (Light)",
         ["Fondo claro base-200", "Texto oscuro legible", "Ideal para ambientes luminosos",
          "Componentes DaisyUI claros"],
         bg_color=RGBColor(240, 245, 250), title_color=ACCENT_BLUE)

# Dark theme card
add_card(slide, Inches(4), Inches(2.1), Inches(2.8), Inches(2.5),
         "🌙  Tema Oscuro (Dracula)",
         ["Fondo oscuro profundo", "Texto claro contrastante", "Reduce fatiga visual",
          "Acentos cyan y púrpura"],
         bg_color=RGBColor(20, 20, 35), title_color=RGBColor(139, 92, 246))

add_multiline_box(slide, Inches(0.8), Inches(4.8), Inches(6), Inches(1),
                  ["La preferencia de tema se guarda en localStorage del navegador.",
                   "Se aplica automáticamente al recargar la página.",
                   "Toggle disponible en la barra de navegación superior."],
                  font_size=11, color=LIGHT_GRAY)

# Sidebar section (right)
add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
             "📌 Sidebar de Navegación", font_size=20, color=ACCENT_CYAN, bold=True)

# Sidebar mockup
add_sidebar_mockup(slide, Inches(7.5), Inches(2.1), Inches(2.3), Inches(4.2),
                   active_item="/dashboard")

# Sidebar features
add_multiline_box(slide, Inches(10.2), Inches(2.1), Inches(2.8), Inches(4.2),
                  ["Características:",
                   "",
                   "✅ Colapsable (toggle)",
                   "✅ Indicador de ruta activa",
                   "   con resaltado cyan",
                   "✅ Íconos SVG por módulo",
                   "✅ Información de usuario",
                   "   en la parte inferior",
                   "✅ Responsive con drawer",
                   "   para pantallas grandes"],
                  font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 - SEGURIDAD
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Seguridad del Sistema", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

security_items = [
    ("🔐 Autenticación JWT", "Tokens firmados con HS256, expiración 60 min.\nCada petición lleva Bearer token en header."),
    ("🔒 Hashing bcrypt", "Las contraseñas se almacenan hasheadas.\nNunca en texto plano en la base de datos."),
    ("🌐 CORS Restringido", "Solo dominios configurados pueden acceder.\nCredentials habilitado para cookies seguras."),
    ("👤 Aislamiento de Datos", "Cada usuario solo ve sus propios datos.\nFiltrado por user_id en cada consulta."),
    ("🛡️ Validación de Entrada", "Pydantic valida toda entrada con tipos estrictos.\nArchivos validados por formato y tamaño."),
    ("🔑 API Keys en .env", "Ninguna clave hardcodeada en el código.\nVariables de entorno excluidas de Git."),
    ("🚪 Guards de Ruta", "AuthGuard protege rutas autenticadas.\nGuestGuard redirige usuarios logueados."),
    ("⏱️ Interceptor HTTP", "Inyecta token automáticamente.\n401 → logout y redirect a login."),
]

for i, (title, desc) in enumerate(security_items):
    row = i // 4
    col = i % 4
    x = Inches(0.5) + col * Inches(3.15)
    y = Inches(1.5) + row * Inches(2.6)
    add_card(slide, x, y, Inches(3), Inches(2.3), title, desc.split("\n"),
             bg_color=RGBColor(20, 30, 50))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 - CIERRE
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Resumen de Funcionalidades", font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(2.6), Inches(3.333))

summary_items = [
    ("🔐", "Autenticación\nJWT Segura"),
    ("📋", "Gestión de\nCampañas"),
    ("🔊", "Carga y Gestión\nde Audios"),
    ("🎙️", "Transcripción\nAutomática"),
    ("🤖", "Análisis QA\ncon IA"),
    ("📊", "Dashboard\nEstadísticas"),
    ("📥", "Exportación\nExcel"),
    ("🎨", "Personalización\nTemas"),
]

for i, (icon, label) in enumerate(summary_items):
    x = Inches(0.5) + i * Inches(1.58)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(3.2),
                                     Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(20, 40, 60)
    circle.line.color.rgb = ACCENT_CYAN
    circle.line.width = Pt(2)
    tf = circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x, Inches(4.3), Inches(1.5), Inches(0.8),
                 label, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "T&A Hub  —  Transcription & Analysis Platform",
             font_size=16, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             "Sistema desarrollado con Angular 21 + FastAPI + PostgreSQL + OpenAI",
             font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             "Abril 2026  •  Versión 1.0",
             font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

output_path = os.path.join(
    r"C:\Users\Luis\Documents\.repositorios_github",
    "TUTORIAL_FUNCIONALIDADES_TA_HUB.pptx"
)
prs.save(output_path)
print(f"Presentación generada exitosamente: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
